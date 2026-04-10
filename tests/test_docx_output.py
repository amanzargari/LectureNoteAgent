from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from lecture_note_agent.docx_utils import write_docx_from_markdown
from lecture_note_agent.io_utils import SlideImageAsset, extract_slide_images


def test_write_docx_from_markdown_embeds_images(tmp_path: Path) -> None:
    image_path = tmp_path / "sample.png"
    Image.new("RGB", (120, 80), color=(255, 0, 0)).save(image_path)

    output_docx = tmp_path / "notes.docx"
    write_docx_from_markdown(
        markdown_text="# Lecture\n\nSee Figure 1.\n\n![Example image](image_ref:img_1)",
        output_path=str(output_docx),
        course_name="Test Course",
        slide_images={
            1: [
                SlideImageAsset(
                    slide_number=1,
                    image_ref="img_1",
                    image_path=str(image_path),
                )
            ]
        },
    )

    assert output_docx.exists()
    with ZipFile(output_docx, "r") as archive:
        media_files = [name for name in archive.namelist() if name.startswith("word/media/")]
    assert media_files, "Expected embedded media files in generated DOCX"


def test_extract_slide_images_applies_pptx_crop(tmp_path: Path) -> None:
    source_image = tmp_path / "src.png"
    Image.new("RGB", (100, 50), color=(0, 128, 255)).save(source_image)

    pptx_path = tmp_path / "deck.pptx"
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    picture = slide.shapes.add_picture(str(source_image), Inches(1), Inches(1), width=Inches(4))
    picture.crop_left = 0.1
    picture.crop_right = 0.1
    prs.save(str(pptx_path))

    images = extract_slide_images(str(pptx_path), artifacts_dir=str(tmp_path / "artifacts"))
    assert 1 in images
    assert images[1]

    extracted_path = Path(images[1][0].image_path)
    assert extracted_path.exists()

    with Image.open(extracted_path) as cropped:
        cropped_width, _ = cropped.size

    assert cropped_width < 100


def test_write_docx_skips_unrecognized_images_without_crashing(tmp_path: Path) -> None:
    bad_img_path = tmp_path / "bad.bin"
    bad_img_path.write_bytes(b"this-is-not-an-image")

    output_docx = tmp_path / "notes_bad_image.docx"
    write_docx_from_markdown(
        markdown_text="# Lecture\n\nBody",
        output_path=str(output_docx),
        course_name="Test Course",
        slide_images={
            1: [
                SlideImageAsset(
                    slide_number=1,
                    image_ref="corrupt_ref",
                    image_path=str(bad_img_path),
                )
            ]
        },
    )

    assert output_docx.exists()


def test_write_docx_embeds_only_referenced_inline_images(tmp_path: Path) -> None:
    used_image = tmp_path / "used.png"
    unused_image = tmp_path / "unused.png"
    Image.new("RGB", (200, 120), color=(255, 200, 100)).save(used_image)
    Image.new("RGB", (80, 80), color=(20, 40, 200)).save(unused_image)

    output_docx = tmp_path / "inline_only.docx"
    write_docx_from_markdown(
        markdown_text=(
            "# Topic\n"
            "See Figure 1 for the architecture overview.\n\n"
            "![Architecture overview](image_ref:used_ref)\n"
        ),
        output_path=str(output_docx),
        course_name="Course",
        slide_images={
            1: [SlideImageAsset(slide_number=1, image_ref="used_ref", image_path=str(used_image))],
            2: [SlideImageAsset(slide_number=2, image_ref="unused_ref", image_path=str(unused_image))],
        },
    )

    assert output_docx.exists()
    with ZipFile(output_docx, "r") as archive:
        media_files = [name for name in archive.namelist() if name.startswith("word/media/")]
        document_xml = archive.read("word/document.xml").decode("utf-8", errors="ignore")

    assert len(media_files) == 1
    assert "Figure 1" in document_xml


def test_write_docx_deduplicates_figure_prefix_in_caption(tmp_path: Path) -> None:
    used_image = tmp_path / "used2.png"
    Image.new("RGB", (200, 120), color=(100, 220, 120)).save(used_image)

    output_docx = tmp_path / "dedupe_figure.docx"
    write_docx_from_markdown(
        markdown_text="![Figure 6: Datacenter Network Architecture](image_ref:used_ref)",
        output_path=str(output_docx),
        course_name="Course",
        slide_images={
            1: [SlideImageAsset(slide_number=1, image_ref="used_ref", image_path=str(used_image))],
        },
    )

    with ZipFile(output_docx, "r") as archive:
        document_xml = archive.read("word/document.xml").decode("utf-8", errors="ignore")

    assert "Figure 1: Datacenter Network Architecture" in document_xml
    assert "Figure 1: Figure 6:" not in document_xml


def test_write_docx_fuzzy_image_ref_matching(tmp_path: Path) -> None:
    used_image = tmp_path / "used3.png"
    Image.new("RGB", (160, 90), color=(80, 160, 220)).save(used_image)

    output_docx = tmp_path / "fuzzy_ref.docx"
    write_docx_from_markdown(
        markdown_text="![Architecture](image_ref:PPTX_SLIDE_1_IMAGE_1.PNG)",
        output_path=str(output_docx),
        course_name="Course",
        slide_images={
            1: [
                SlideImageAsset(
                    slide_number=1,
                    image_ref="pptx_slide_1_image_1",
                    image_path=str(used_image),
                )
            ],
        },
    )

    with ZipFile(output_docx, "r") as archive:
        media_files = [name for name in archive.namelist() if name.startswith("word/media/")]

    assert len(media_files) == 1
