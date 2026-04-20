from __future__ import annotations

import hashlib
import os
import re
from io import BytesIO
from dataclasses import dataclass
from pathlib import Path
from typing import List

from pypdf import PdfReader
from pptx import Presentation

from .models import SlideUnit, TranscriptSegment


MATH_PATTERNS = [
    r"\b\d+\s*[+\-*/=]\s*\d+\b",
    r"\b(?:sin|cos|tan|log|ln|exp|lim|sum|prod|int)\s*\(",
    r"\b[A-Za-z]\s*=\s*[^\n]{1,40}",
    r"\$[^$]+\$",
]


@dataclass(frozen=True)
class SlideImageAsset:
    slide_number: int
    image_ref: str
    image_path: str


def _env_int(name: str, default: int, *, minimum: int = 0) -> int:
    raw = os.getenv(name)
    if raw is None:
        return max(minimum, default)

    try:
        value = int(raw.strip())
    except (TypeError, ValueError):
        return max(minimum, default)

    return max(minimum, value)


def _env_float(name: str, default: float, *, minimum: float = 0.0) -> float:
    raw = os.getenv(name)
    if raw is None:
        return max(minimum, default)

    try:
        value = float(raw.strip())
    except (TypeError, ValueError):
        return max(minimum, default)

    return max(minimum, value)


def _pdf_image_filter_config() -> tuple[int, int, float, int]:
    """Runtime-tunable guardrails for filtering noisy PDF image objects."""
    min_area = _env_int("PDF_IMAGE_MIN_AREA", 7000, minimum=1)
    min_edge = _env_int("PDF_IMAGE_MIN_EDGE", 45, minimum=1)
    max_aspect_ratio = _env_float("PDF_IMAGE_MAX_ASPECT_RATIO", 8.0, minimum=1.0)
    max_per_slide = _env_int("PDF_MAX_IMAGES_PER_SLIDE", 3, minimum=1)
    return min_area, min_edge, max_aspect_ratio, max_per_slide


def _pdf_image_dimensions(image_obj: object) -> tuple[int, int]:
    pil_img = getattr(image_obj, "image", None)
    if pil_img is not None:
        try:
            width_px, height_px = pil_img.size
            return int(width_px), int(height_px)
        except Exception:
            pass

    width = getattr(image_obj, "width", None)
    height = getattr(image_obj, "height", None)
    if isinstance(width, int) and isinstance(height, int) and width > 0 and height > 0:
        return width, height

    blob = getattr(image_obj, "data", None)
    if isinstance(blob, (bytes, bytearray)) and blob:
        try:
            from PIL import Image

            with Image.open(BytesIO(bytes(blob))) as decoded:
                width_px, height_px = decoded.size
                return int(width_px), int(height_px)
        except Exception:
            return 0, 0

    return 0, 0


def _pdf_image_hash(image_obj: object) -> str:
    blob = getattr(image_obj, "data", None)
    if isinstance(blob, (bytes, bytearray)) and blob:
        return hashlib.sha1(bytes(blob)).hexdigest()

    pil_img = getattr(image_obj, "image", None)
    if pil_img is not None:
        try:
            return hashlib.sha1(pil_img.tobytes()).hexdigest()
        except Exception:
            return ""

    return ""


def _is_high_signal_pdf_image(
    *,
    width_px: int,
    height_px: int,
    min_area: int,
    min_edge: int,
    max_aspect_ratio: float,
) -> bool:
    if width_px <= 0 or height_px <= 0:
        return False

    if width_px * height_px < min_area:
        return False

    if min(width_px, height_px) < min_edge:
        return False

    aspect_ratio = max(width_px, height_px) / max(1, min(width_px, height_px))
    if aspect_ratio > max_aspect_ratio:
        return False

    return True


def _select_pdf_images_for_page(page: object, slide_number: int) -> list[tuple[int, str, object]]:
    images = list(getattr(page, "images", []) or [])
    if not images:
        return []

    min_area, min_edge, max_aspect_ratio, max_per_slide = _pdf_image_filter_config()
    seen_hashes: set[str] = set()
    candidates: list[tuple[int, int, str, object]] = []

    for image_idx, image_obj in enumerate(images, 1):
        image_ref = getattr(image_obj, "name", "") or f"pdf_page_{slide_number}_image_{image_idx}"
        width_px, height_px = _pdf_image_dimensions(image_obj)
        if not _is_high_signal_pdf_image(
            width_px=width_px,
            height_px=height_px,
            min_area=min_area,
            min_edge=min_edge,
            max_aspect_ratio=max_aspect_ratio,
        ):
            continue

        image_hash = _pdf_image_hash(image_obj)
        if image_hash and image_hash in seen_hashes:
            continue
        if image_hash:
            seen_hashes.add(image_hash)

        area = width_px * height_px
        candidates.append((area, image_idx, image_ref, image_obj))

    if not candidates:
        return []

    # Prefer larger, information-dense visuals; keep ordering stable for readability.
    candidates.sort(key=lambda item: (-item[0], item[1]))
    selected = candidates[:max_per_slide]
    selected.sort(key=lambda item: item[1])

    return [(image_idx, image_ref, image_obj) for _, image_idx, image_ref, image_obj in selected]


def _extract_formula_candidates(text: str) -> List[str]:
    matches: list[str] = []
    for pattern in MATH_PATTERNS:
        for m in re.finditer(pattern, text):
            token = m.group(0).strip()
            if token and token not in matches:
                matches.append(token)
    return matches[:30]


def extract_formula_candidates(text: str) -> List[str]:
    return _extract_formula_candidates(text)


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _has_meaningful_text(text: str, min_alnum: int = 30) -> bool:
    alnum_count = sum(ch.isalnum() for ch in text)
    return alnum_count >= min_alnum


def has_meaningful_text(text: str, min_alnum: int = 30) -> bool:
    return _has_meaningful_text(text, min_alnum=min_alnum)


def parse_slides(
    slides_path: str,
) -> List[SlideUnit]:
    path = Path(slides_path)
    ext = path.suffix.lower()

    if ext in {".md", ".txt"}:
        text = _read_text_file(path)
        chunks = [c.strip() for c in re.split(r"\n---+\n|\n#{1,2}\s+Slide", text) if c.strip()]
        slides: list[SlideUnit] = []
        for idx, chunk in enumerate(chunks, 1):
            lines = [ln.strip() for ln in chunk.splitlines() if ln.strip()]
            title = lines[0][:140] if lines else f"Slide {idx}"
            body = "\n".join(lines[1:]) if len(lines) > 1 else chunk
            slides.append(
                SlideUnit(
                    slide_number=idx,
                    title=title,
                    text=body,
                    image_refs=[],
                    formula_candidates=_extract_formula_candidates(chunk),
                )
            )
        return slides

    if ext == ".pdf":
        reader = PdfReader(str(path))
        slides = []
        for i, page in enumerate(reader.pages, 1):
            page_text = (page.extract_text() or "").strip()
            full_text = page_text
            selected_images = _select_pdf_images_for_page(page, i)
            image_refs = [image_ref for _, image_ref, _ in selected_images]
            title = full_text.splitlines()[0][:140] if full_text else f"Slide {i}"
            slides.append(
                SlideUnit(
                    slide_number=i,
                    title=title,
                    text=full_text,
                    image_refs=image_refs,
                    formula_candidates=_extract_formula_candidates(full_text),
                )
            )
        return slides

    if ext == ".pptx":
        prs = Presentation(str(path))
        slides: list[SlideUnit] = []
        for i, slide in enumerate(prs.slides, 1):
            text_parts: list[str] = []
            image_refs: list[str] = []
            title = ""
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text:
                    shape_text = shape.text.strip()
                    if shape_text:
                        text_parts.append(shape_text)
                        if not title:
                            title = shape_text.splitlines()[0][:140]
                if getattr(shape, "shape_type", None) == 13:  # Picture
                    shape_name = getattr(shape, "name", "") or f"pptx_slide_{i}_image_{len(image_refs)+1}"
                    image_refs.append(shape_name)
            text = "\n".join(text_parts).strip()
            slides.append(
                SlideUnit(
                    slide_number=i,
                    title=title or f"Slide {i}",
                    text=text,
                    image_refs=image_refs,
                    formula_candidates=_extract_formula_candidates(text),
                )
            )
        return slides

    raise ValueError(f"Unsupported slides format: {ext}. Use .pdf, .pptx, .md, or .txt")


_SRT_TIMESTAMP = re.compile(
    r"(\d{1,2}:\d{2}:\d{2})[,.](\d{1,3})\s*-->\s*\d{1,2}:\d{2}:\d{2}[,.]\d{1,3}"
)
_SRT_INDEX = re.compile(r"^\d+$")


def _parse_srt(text: str) -> List[TranscriptSegment]:
    segments: list[TranscriptSegment] = []
    seg_idx = 0
    lines = text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        # Skip sequence number lines
        if _SRT_INDEX.match(line):
            i += 1
            continue
        ts_match = _SRT_TIMESTAMP.match(line)
        if ts_match:
            timestamp = ts_match.group(1)  # HH:MM:SS
            i += 1
            text_lines: list[str] = []
            while i < len(lines) and lines[i].strip():
                text_lines.append(lines[i].strip())
                i += 1
            content = " ".join(text_lines).strip()
            if content:
                seg_idx += 1
                segments.append(
                    TranscriptSegment(
                        segment_id=f"T{seg_idx}",
                        timestamp=timestamp,
                        speaker="Teacher",
                        text=content,
                    )
                )
        else:
            i += 1
    return segments


def parse_transcript(transcript_path: str) -> List[TranscriptSegment]:
    path = Path(transcript_path)
    text = _read_text_file(path)

    if path.suffix.lower() == ".srt" or _SRT_TIMESTAMP.search(text):
        return _parse_srt(text)

    lines = [ln.rstrip() for ln in text.splitlines() if ln.strip()]
    segments: list[TranscriptSegment] = []

    ts_pattern = re.compile(r"^\[?(\d{1,2}:\d{2}(?::\d{2})?)\]?\s*(.*)$")
    speaker_pattern = re.compile(r"^([A-Za-z][A-Za-z0-9_\- ]{1,30}):\s*(.*)$")

    for i, line in enumerate(lines, 1):
        timestamp = ""
        speaker = "Teacher"
        content = line

        ts_match = ts_pattern.match(content)
        if ts_match:
            timestamp = ts_match.group(1)
            content = ts_match.group(2).strip()

        sp_match = speaker_pattern.match(content)
        if sp_match:
            speaker = sp_match.group(1).strip()
            content = sp_match.group(2).strip()

        segments.append(
            TranscriptSegment(
                segment_id=f"T{i}",
                timestamp=timestamp,
                speaker=speaker,
                text=content,
            )
        )

    return segments


def build_source_payload(course_name: str, slides: List[SlideUnit], transcript: List[TranscriptSegment]) -> str:
    slide_blocks: list[str] = []
    for s in slides:
        img_line = ", ".join(s.image_refs) if s.image_refs else "None"
        formula_line = ", ".join(s.formula_candidates) if s.formula_candidates else "None"
        slide_blocks.append(
            "\n".join(
                [
                    f"[S{s.slide_number}] Title: {s.title}",
                    f"Text: {s.text[:3000]}",
                    f"ImageRefs: {img_line}",
                    f"FormulaCandidates: {formula_line}",
                ]
            )
        )

    transcript_blocks = [
        f"[{t.segment_id}] ({t.timestamp or 'NA'}) {t.speaker}: {t.text}"
        for t in transcript
    ]

    payload = (
        f"Course: {course_name}\n\n"
        f"## Slides ({len(slides)})\n" + "\n\n".join(slide_blocks) + "\n\n"
        f"## Transcript Segments ({len(transcript)})\n" + "\n".join(transcript_blocks)
    )
    return payload


def _normalize_crop_fraction(value: object) -> float:
    if value is None:
        return 0.0
    try:
        out = float(value)
    except (TypeError, ValueError):
        return 0.0

    if out > 1.0:
        out = out / 100000.0
    if out < 0:
        return 0.0
    return min(out, 0.95)


def _crop_image(
    source_path: Path,
    output_path: Path,
    crop_left: float,
    crop_top: float,
    crop_right: float,
    crop_bottom: float,
) -> Path:
    from PIL import Image

    relax_factor = max(0.0, min(1.0, float(os.getenv("IMAGE_CROP_RELAX_FACTOR", "0.75"))))
    min_retain = max(0.25, min(0.95, float(os.getenv("IMAGE_CROP_MIN_RETAIN_RATIO", "0.45"))))

    eff_left = crop_left * relax_factor
    eff_top = crop_top * relax_factor
    eff_right = crop_right * relax_factor
    eff_bottom = crop_bottom * relax_factor

    with Image.open(source_path) as img:
        width, height = img.size
        left = int(width * eff_left)
        top = int(height * eff_top)
        right = int(width * (1.0 - eff_right))
        bottom = int(height * (1.0 - eff_bottom))

        if right <= left or bottom <= top:
            img.save(output_path, format="PNG")
            return output_path

        retained_w_ratio = (right - left) / max(1, width)
        retained_h_ratio = (bottom - top) / max(1, height)
        retained_area_ratio = retained_w_ratio * retained_h_ratio

        # If crop appears too aggressive, keep full image for readability.
        if (
            retained_w_ratio < min_retain
            or retained_h_ratio < min_retain
            or retained_area_ratio < (min_retain * min_retain)
        ):
            img.save(output_path, format="PNG")
            return output_path

        cropped = img.crop((left, top, right, bottom))
        cropped.save(output_path, format="PNG")
    return output_path


def _safe_suffix(suffix: str, default: str = ".png") -> str:
    normalized = (suffix or "").strip().lower()
    if normalized.startswith("."):
        ext = normalized
    else:
        ext = f".{normalized}" if normalized else default

    if ext in {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"}:
        return ext
    return default


def extract_slide_images(slides_path: str, artifacts_dir: str | None = None) -> dict[int, list[SlideImageAsset]]:
    """Extract concrete slide images for embedding in final DOCX output.

    For PPTX, respects picture crop metadata and exports cropped images.
    """
    source = Path(slides_path)
    ext = source.suffix.lower()

    if artifacts_dir:
        root = Path(artifacts_dir)
    else:
        root = source.parent
    image_root = root / "extracted_images"
    image_root.mkdir(parents=True, exist_ok=True)

    slide_images: dict[int, list[SlideImageAsset]] = {}

    if ext == ".pdf":
        reader = PdfReader(str(source))
        for slide_number, page in enumerate(reader.pages, 1):
            assets: list[SlideImageAsset] = []
            selected_images = _select_pdf_images_for_page(page, slide_number)
            for _, image_ref, img in selected_images:
                export_idx = len(assets) + 1
                pil_img = getattr(img, "image", None)
                if pil_img is not None:
                    png_path = image_root / f"slide_{slide_number}_{export_idx}.png"
                    pil_img.save(png_path)
                    assets.append(
                        SlideImageAsset(
                            slide_number=slide_number,
                            image_ref=image_ref,
                            image_path=str(png_path),
                        )
                    )
                    continue

                suffix = _safe_suffix(Path(image_ref).suffix)
                image_path = image_root / f"slide_{slide_number}_{export_idx}{suffix}"
                blob = getattr(img, "data", None)
                if isinstance(blob, (bytes, bytearray)) and blob:
                    image_path.write_bytes(bytes(blob))
                    assets.append(
                        SlideImageAsset(
                            slide_number=slide_number,
                            image_ref=image_ref,
                            image_path=str(image_path),
                        )
                    )

            if assets:
                slide_images[slide_number] = assets

        return slide_images

    if ext == ".pptx":
        prs = Presentation(str(source))
        for slide_number, slide in enumerate(prs.slides, 1):
            assets: list[SlideImageAsset] = []
            image_idx = 0

            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) != 13:
                    continue

                image_idx += 1
                image_ref = getattr(shape, "name", "") or f"pptx_slide_{slide_number}_image_{image_idx}"
                pic = getattr(shape, "image", None)
                if pic is None:
                    continue

                suffix = _safe_suffix(getattr(pic, "ext", "png"))
                original_path = image_root / f"slide_{slide_number}_{image_idx}_orig{suffix}"
                original_path.write_bytes(pic.blob)

                crop_left = _normalize_crop_fraction(getattr(shape, "crop_left", 0.0))
                crop_top = _normalize_crop_fraction(getattr(shape, "crop_top", 0.0))
                crop_right = _normalize_crop_fraction(getattr(shape, "crop_right", 0.0))
                crop_bottom = _normalize_crop_fraction(getattr(shape, "crop_bottom", 0.0))

                if any(v > 0 for v in (crop_left, crop_top, crop_right, crop_bottom)):
                    cropped_path = image_root / f"slide_{slide_number}_{image_idx}_cropped.png"
                    final_path = _crop_image(
                        source_path=original_path,
                        output_path=cropped_path,
                        crop_left=crop_left,
                        crop_top=crop_top,
                        crop_right=crop_right,
                        crop_bottom=crop_bottom,
                    )
                else:
                    final_path = original_path

                assets.append(
                    SlideImageAsset(
                        slide_number=slide_number,
                        image_ref=image_ref,
                        image_path=str(final_path),
                    )
                )

            if assets:
                slide_images[slide_number] = assets

        return slide_images

    return slide_images
